from pptx import Presentation
from pptx.util import Inches


def img2pptx(img_paths, pptx_path, aspect_ratio=(4, 3)):
    presentation = Presentation()

    slide_width = Inches(aspect_ratio[0])
    slide_height = Inches(aspect_ratio[1])
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    for img_path in img_paths:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        slide.shapes.add_picture(img_path, 0, 0,  width=slide_width, height=slide_height)

    presentation.save(pptx_path)