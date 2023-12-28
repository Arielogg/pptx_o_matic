import fitz
import io
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import os


def pdf_to_images(pdf_path):
    src = fitz.open(pdf_path)
    mat = fitz.Matrix(7,7)
    images = []
    for page in src:
        pix = page.get_pixmap(matrix=mat)
        img = Image.open(io.BytesIO(pix.tobytes('png')))
        images.append(img)
    return images


def pdf2pptx(pdf_path, pptx_path, aspect_ratio=(4, 3)):
    images = pdf_to_images(pdf_path)
    presentation = Presentation()

    # Set slide width and height according to the aspect ratio
    slide_width = Inches(aspect_ratio[0])
    slide_height = Inches(aspect_ratio[1])
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    for image in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        image_path = 'temp_img.jpg'
        image.save(image_path, 'JPEG')

        slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)
        os.remove(image_path)

    presentation.save(pptx_path)